from typing import Any, Dict, List, Tuple, Optional
import os, io, json, re
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import fitz

# Optional EasyOCR (only used for bitmap PDFs/images)
try:
    import easyocr  # type: ignore
    _easyocr_reader = None
except Exception:
    easyocr = None
    _easyocr_reader = None


# =========================
# Extraction helpers
# =========================
def _extract_text_from_pptx(pptx_path: Path) -> str:
    prs = Presentation(pptx_path)
    chunks: List[str] = []
    for slide in prs.slides:
        # shapes text
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
            except Exception:
                continue
        # notes
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                nt = slide.notes_slide.notes_text_frame.text
                if nt:
                    chunks.append(nt)
        except Exception:
            pass
    return "\n\n".join(chunks)


def _extract_text_from_pdf(pdf_path: Path, use_ocr: bool = False, lang: str = "en") -> str:
    doc = fitz.open(pdf_path)
    parts: List[str] = []
    has_text = False
    for page in doc:
        txt = page.get_text().strip()
        if txt:
            has_text = True
            parts.append(txt)
    if has_text or not use_ocr:
        return "\n\n".join(parts)

    # OCR path
    if easyocr is None:
        return "\n\n".join(parts)

    global _easyocr_reader
    if _easyocr_reader is None:
        _easyocr_reader = easyocr.Reader([lang], gpu=False)

    oc_parts: List[str] = []
    for page in doc:
        pix = page.get_pixmap(dpi=220)
        import numpy as np, cv2  # local import
        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
        if img.shape[2] == 4:
            img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)
        results = _easyocr_reader.readtext(img, detail=0, paragraph=True)
        if results:
            oc_parts.append("\n".join(results))
    return "\n\n".join(oc_parts)


# =========================
# Gemini helpers (inline)
# =========================
def _get_api_key(explicit: Optional[str]) -> Optional[str]:
    return explicit or os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")


def _gemini_generate_qa(
    text: str,
    total: int = 20,
    model_name: str = "gemini-1.5-flash",
    api_key: Optional[str] = None
) -> List[Tuple[str, str]]:
    """
    Simple Q/A generator (question, answer) using Gemini, with robust parsing.
    """
    key = _get_api_key(api_key)
    if not key:
        raise RuntimeError("Gemini API key missing. Set GEMINI_API_KEY or pass gemini_api_key.")

    import google.generativeai as genai
    genai.configure(api_key=key)
    model = genai.GenerativeModel(model_name)
    sys = (
        "You are a question generator. Read the provided lecture text and produce concise study questions "
        "with their short, correct answers. Return STRICT JSON: a list of objects with keys 'q' and 'a'. "
        "Avoid markdown or commentary."
    )
    prompt = f"{sys}\n\nTarget count: {total}\nLecture text:\n{text[:12000]}"

    resp = model.generate_content(prompt)
    out: List[Tuple[str, str]] = []
    try:
        data = getattr(resp, "text", "") or ""
        # strip code fences / grab first [...]
        data = re.sub(r"^```(?:json)?\s*|\s*```$", "", data.strip(), flags=re.I | re.M)
        if "[" not in data:
            # try candidate path
            if getattr(resp, "candidates", None):
                data = resp.candidates[0].content.parts[0].text  # type: ignore
        start = data.find('['); end = data.rfind(']')
        if start != -1 and end != -1:
            data = data[start:end+1]
        arr = json.loads(data)
        for item in arr:
            q = (item.get("q") or item.get("question") or "").strip()
            a = (item.get("a") or item.get("answer") or "").strip()
            if q and a:
                out.append((q, a))
    except Exception:
        # Q:/A: fallback parse
        text_resp = getattr(resp, "text", "") or ""
        lines = [ln.strip() for ln in text_resp.splitlines() if ln.strip()]
        q, a = None, None
        for ln in lines:
            if ln.lower().startswith("q:"):
                if q and a:
                    out.append((q, a))
                q, a = ln[2:].strip(), ""
            elif ln.lower().startswith("a:") and q is not None:
                a = ln[2:].strip()
        if q and a:
            out.append((q, a))
    return out[:total]


def _gemini_explain_batch(
    qa: List[Tuple[str, str]],
    text: str,
    model_name: str,
    api_key: Optional[str]
) -> List[str]:
    """
    2–3 sentence explanation per QA (same order). Robust JSON parsing.
    """
    key = _get_api_key(api_key)
    if not key:
        # If you want to make explanations optional, return heuristics instead of raising:
        return [f"Explanation: {a[:200]}" for _, a in qa]

    import google.generativeai as genai
    genai.configure(api_key=key)
    model = genai.GenerativeModel(model_name)

    pack = [{"q": q, "a": a} for q, a in qa]
    prompt = (
        "For each item, produce a short 2–3 sentence explanation (no formatting). "
        "Return STRICT JSON: a list of strings in the same order; its length must equal the input list.\n\n"
        f"LECTURE (truncated):\n{text[:8000]}\n\n"
        f"ITEMS:\n{json.dumps(pack, ensure_ascii=False)}"
    )

    resp = model.generate_content(prompt)
    data = getattr(resp, "text", "") or "[]"
    data = re.sub(r"^```(?:json)?\s*|\s*```$", "", data.strip(), flags=re.I | re.M)
    start, end = data.find("["), data.rfind("]")
    if start != -1 and end != -1:
        data = data[start:end+1]
    try:
        arr = json.loads(data)
        out = [str(x) for x in arr]
        if len(out) != len(qa):
            return [f"Explanation: {a[:200]}" for _, a in qa]
        return out
    except Exception:
        return [f"Explanation: {a[:200]}" for _, a in qa]


def _infer_title(
    lecture_text: str,
    filename_stem: str,
    model_name: str,
    api_key: Optional[str]
) -> str:
    """
    Try Gemini for a concise 3–7 word title; otherwise heuristic or file name.
    """
    key = _get_api_key(api_key)
    if key:
        try:
            import google.generativeai as genai
            genai.configure(api_key=key)
            m = genai.GenerativeModel(model_name)
            prompt = (
                "Give a concise 3–7 word title for this lecture text. "
                "Return ONLY the title, no quotes, no trailing punctuation.\n\n"
                f"{lecture_text[:8000]}"
            )
            t = (m.generate_content(prompt).text or "").replace("\n", " ").strip(" .,:;\"'")
            if len(t.split()) >= 2:
                return t[:80]
        except Exception:
            pass
    # Heuristic fallback
    for line in lecture_text.splitlines():
        s = re.sub(r"[^A-Za-z0-9 :/\-\(\)\[\]]+", "", line).strip()
        if len(s.split()) >= 2 and len(s) >= 10:
            return s[:80]
    return filename_stem[:80] or "Auto Quiz"


# =========================
# Slide builder (improved)
# =========================
ACCENT_Q = RGBColor(37, 99, 235)    # blue
ACCENT_A = RGBColor(16, 185, 129)   # green
ACCENT_X = RGBColor(234, 179, 8)    # amber
FONT_NAME = "Calibri"

def _set_text_frame(tf, text: str, size_pt: int = 24, bold: bool = False,
                    color: Optional[RGBColor] = None, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    for para in tf.paragraphs:
        para.space_after = Pt(6)
        para.space_before = Pt(0)
    tf.word_wrap = True


def _build_qa_deck(
    qa: List[Tuple[str, str]],
    title: str = "Auto Quiz",
    explanations: Optional[List[str]] = None,
    source_name: str = ""
) -> bytes:
    prs = Presentation()

    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    try:
        _set_text_frame(s0.shapes.title.text_frame, title, 44, True)
    except Exception:
        # if template differs, fallback to textbox
        tx = s0.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1.5)).text_frame
        _set_text_frame(tx, title, 44, True)
    if len(s0.placeholders) > 1:
        sub = f"Generated quiz deck • {date.today().isoformat()}"
        if source_name:
            sub += f" • from {source_name}"
        try:
            _set_text_frame(s0.placeholders[1].text_frame, sub, 18)
        except Exception:
            pass

    total = len(qa)
    for i, (q, a) in enumerate(qa, start=1):
        # Question
        sq = prs.slides.add_slide(prs.slide_layouts[1])  # Title+Content
        try:
            _set_text_frame(sq.shapes.title.text_frame, f"Question {i}/{total}", 30, True, ACCENT_Q)
            _set_text_frame(sq.placeholders[1].text_frame, q, 28)
        except Exception:
            # fallback to textboxes
            _set_textbox(sq, f"Question {i}/{total}", title=True, color=ACCENT_Q)
            _set_textbox(sq, q, top_inches=2.2)

        # Answer
        sa = prs.slides.add_slide(prs.slide_layouts[1])
        try:
            _set_text_frame(sa.shapes.title.text_frame, f"Answer {i}/{total}", 30, True, ACCENT_A)
            _set_text_frame(sa.placeholders[1].text_frame, a, 28)
        except Exception:
            _set_textbox(sa, f"Answer {i}/{total}", title=True, color=ACCENT_A)
            _set_textbox(sa, a, top_inches=2.2)

        # Explanation (optional)
        if explanations:
            sx = prs.slides.add_slide(prs.slide_layouts[1])
            text_x = explanations[i - 1]
            try:
                _set_text_frame(sx.shapes.title.text_frame, f"Why? {i}/{total}", 30, True, ACCENT_X)
                _set_text_frame(sx.placeholders[1].text_frame, text_x, 24)
            except Exception:
                _set_textbox(sx, f"Why? {i}/{total}", title=True, color=ACCENT_X)
                _set_textbox(sx, text_x, top_inches=2.2, size=24)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _set_textbox(slide, text: str, title: bool = False,
                 color: Optional[RGBColor] = None, top_inches: float = 1.5, size: int = 24):
    left, top, width, height = Inches(1), Inches(top_inches), Inches(8), Inches(4.5)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28 if title else size)
    if color:
        run.font.color.rgb = color
    if title:
        run.font.bold = True


# =========================
# Orchestration
# =========================
def run_pipeline_end_to_end(
    pptx_file,
    ocr_engine: str = "easyocr",
    language: str = "en",
    prefer_com: bool = False,
    dpi: int = 180,
    gemini_api_key: Optional[str] = None,
    model_name: str = "gemini-1.5-flash",
    total_questions: int = 20,
    mix_mode: str = "balanced",
    mcq_n: int = 0,
    theory_n: int = 0,
    codefill_n: int = 0,
    difficulty: str = "mixed",
    deck_title: str = "Auto Quiz",
    include_thumbs: bool = True   # repurposed: include explanations when True
):
    """
    Returns: (pptx_bytes, zip_path, out_dir, msg)
    """
    # --- Save the upload with the correct extension
    tmp_dir = Path(os.getenv("TMP", str(Path.cwd() / "tmp")))
    tmp_dir.mkdir(parents=True, exist_ok=True)

    ext = ""
    data: Optional[bytes] = None
    name_hint = getattr(pptx_file, "name", "") or ""

    if hasattr(pptx_file, "read"):  # file-like
        data = pptx_file.read()
        if name_hint:
            ext = Path(name_hint).suffix.lower()
    else:  # path-like
        p = Path(pptx_file)
        data = p.read_bytes()
        ext = p.suffix.lower()

    if not ext:
        if data and data.startswith(b"%PDF"):
            ext = ".pdf"
        elif data and data[:2] == b"PK":
            ext = ".pptx"

    if ext not in (".pptx", ".pdf"):
        raise ValueError(f"Unsupported file type: {ext or 'unknown'}. Upload .pptx or .pdf")

    tmp_in = tmp_dir / f"input_upload{ext}"
    tmp_in.write_bytes(data or b"")

    # --- Extract text
    if ext == ".pptx":
        lecture_text = _extract_text_from_pptx(tmp_in)
    else:
        lecture_text = _extract_text_from_pdf(tmp_in, use_ocr=(ocr_engine == "easyocr"), lang=language)

    if not lecture_text.strip():
        raise ValueError("No text extracted from the document. Check OCR/inputs.")

    # --- Generate Q/A
    qa = _gemini_generate_qa(lecture_text, total=total_questions, model_name=model_name, api_key=gemini_api_key)
    if not qa:
        raise ValueError("Q/A generation returned empty results.")

    # --- Title (Gemini → heuristic) & Explanations (optional)
    filename_stem = Path(name_hint or tmp_in.name).stem
    auto_title = _infer_title(lecture_text, filename_stem, model_name, gemini_api_key)
    final_title = deck_title if (deck_title and deck_title != "Auto Quiz") else auto_title

    explanations = _gemini_explain_batch(qa, lecture_text, model_name, gemini_api_key) if include_thumbs else None

    # --- Build deck (higher-quality)
    pptx_bytes = _build_qa_deck(qa, title=final_title, explanations=explanations, source_name=filename_stem)
    return pptx_bytes, None, None, "ok"
