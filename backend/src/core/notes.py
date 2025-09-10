# -*- coding: utf-8 -*-
"""
OCR + notes utilities for slide text extraction.

Improvements:
- Robust imports with graceful fallbacks
- Higher-quality OCR pre-processing (denoise/threshold/resize)
- Smarter Markdown: heading detection, bullet/number preservation, indentation from left offsets
- Safer temp/IO handling and predictable return schema
"""

from typing import Any, Dict, List, Tuple, Optional
from pathlib import Path
import io
import os
import re
import zipfile

# Third-party (guarded)
try:
    import numpy as np
except Exception:  # pragma: no cover
    np = None  # type: ignore

try:
    import cv2  # type: ignore
except Exception:  # pragma: no cover
    cv2 = None  # type: ignore

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None  # type: ignore

try:
    from pptx import Presentation  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore

try:
    from tqdm import tqdm  # progress bar
except Exception:  # pragma: no cover
    def tqdm(x, **_: Any):
        return x

# Local, optional helpers
# ocr engines
try:
    from .ocr import get_ocr_engine, ocr_image_easy, ocr_image_paddle  # type: ignore
except Exception:  # pragma: no cover
    get_ocr_engine = None  # type: ignore
    def ocr_image_easy(engine, img_np):  # type: ignore
        raise RuntimeError("ocr_image_easy not available")
    def ocr_image_paddle(engine, img_np):  # type: ignore
        raise RuntimeError("ocr_image_paddle not available")

# slide rendering
try:
    from .render import render_slides_to_images  # type: ignore
except Exception:  # pragma: no cover
    render_slides_to_images = None  # type: ignore

# cache cleanup
try:
    from .helpers import cleanup_ocr_cache  # type: ignore
except Exception:  # pragma: no cover
    def cleanup_ocr_cache(_engine):  # type: ignore
        return None


# -----------------------------
# Notes + Alt text extraction
# -----------------------------
def extract_notes_and_alttext(pptx_path: str) -> Dict[int, Dict[str, Any]]:
    """
    Extract notes and alternative text per slide from a .pptx.

    Returns:
        {slide_index: {"notes": str, "alt_texts": List[str]}}
    """
    out: Dict[int, Dict[str, Any]] = {}
    p = Path(pptx_path)
    if p.suffix.lower() != ".pptx":
        return out
    if Presentation is None:
        return out

    prs = Presentation(str(p))
    for i, slide in enumerate(prs.slides, 1):
        # Notes
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            pass

        # Alt text (on shapes/pictures)
        alt_texts: List[str] = []
        for shape in slide.shapes:
            try:
                alt = getattr(shape, "alternative_text", "") or ""
                if alt and isinstance(alt, str):
                    alt = alt.strip()
                    if alt:
                        alt_texts.append(alt)
            except Exception:
                continue

        out[i] = {"notes": notes, "alt_texts": alt_texts}
    return out


# -----------------------------
# OCR utilities
# -----------------------------
def bbox_to_xywh(b: List[Tuple[float, float]]) -> Tuple[float, float, float, float]:
    xs = [p[0] for p in b]; ys = [p[1] for p in b]
    x, y = min(xs), min(ys)
    w, h = max(xs) - x, max(ys) - y
    return x, y, w, h


def _preprocess_for_ocr(img_np, enable_denoise: bool = True, enable_binarize: bool = True):
    """
    Optional denoise + adaptive threshold for sharper OCR, if OpenCV is available.
    """
    if cv2 is None or np is None:
        return img_np

    work = img_np
    if work.ndim == 3 and work.shape[2] == 4:
        work = cv2.cvtColor(work, cv2.COLOR_BGRA2BGR)
    if work.ndim == 3:
        gray = cv2.cvtColor(work, cv2.COLOR_BGR2GRAY)
    else:
        gray = work

    if enable_denoise:
        gray = cv2.fastNlMeansDenoising(gray, h=7)

    if enable_binarize:
        # adaptive threshold is robust for slides with gradients
        th = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY, 31, 10
        )
        return th

    return gray


def normalize_lines(items: List[Dict[str, Any]], y_merge: int = 12) -> List[Dict[str, Any]]:
    """
    items: list of { 'bbox': [(x,y),...], 'text': str, 'conf': float }
    Groups by rows, sorts left-to-right, merges fragments, returns clean lines with median y/left and average conf.
    """
    if np is None:
        # best-effort without numpy
        enriched = []
        for it in items:
            xywh = bbox_to_xywh(it["bbox"])
            enriched.append({**it, "xywh": xywh})
        enriched.sort(key=lambda d: (round(d["xywh"][1] / y_merge), d["xywh"][0]))

        merged, current, last_row = [], None, None
        for it in enriched:
            x, y, w, h = it["xywh"]
            text = (it.get("text") or "").strip()
            if not text:
                continue
            row = round(y / y_merge)
            if current is None:
                current = {"row": row, "x": x, "texts": [text], "xs": [x], "ys": [y], "conf": [it.get("conf", 1.0)]}
                last_row = row
            else:
                if row == last_row:
                    current["texts"].append(text)
                    current["xs"].append(x); current["ys"].append(y); current["conf"].append(it.get("conf", 1.0))
                else:
                    merged.append(current)
                    current = {"row": row, "x": x, "texts": [text], "xs": [x], "ys": [y], "conf": [it.get("conf", 1.0)]}
                    last_row = row
        if current:
            merged.append(current)

        lines = []
        for m in merged:
            line_text = re.sub(r"\s+", " ", " ".join(m["texts"]).strip())
            if not line_text:
                continue
            avg_conf = sum(m["conf"]) / max(1, len(m["conf"]))
            ys_sorted = sorted(m["ys"])
            mid = ys_sorted[len(ys_sorted) // 2]
            lines.append({"text": line_text, "left": min(m["xs"]), "y": float(mid), "conf": float(avg_conf)})
        return lines

    # numpy implementation
    enriched = []
    for it in items:
        xywh = bbox_to_xywh(it["bbox"])
        enriched.append({**it, "xywh": xywh})
    enriched = sorted(enriched, key=lambda d: (round(d["xywh"][1] / y_merge), d["xywh"][0]))

    merged, current, last_row = [], None, None
    for it in enriched:
        x, y, w, h = it["xywh"]
        text = (it.get("text") or "").strip()
        if not text:
            continue
        row = round(y / y_merge)
        if current is None:
            current = {"row": row, "x": x, "texts": [text], "xs": [x], "ys": [y], "conf": [it.get("conf", 1.0)]}
            last_row = row
        else:
            if row == last_row:
                current["texts"].append(text)
                current["xs"].append(x); current["ys"].append(y); current["conf"].append(it.get("conf", 1.0))
            else:
                merged.append(current)
                current = {"row": row, "x": x, "texts": [text], "xs": [x], "ys": [y], "conf": [it.get("conf", 1.0)]}
                last_row = row
    if current:
        merged.append(current)

    lines = []
    for m in merged:
        # merge fragments; fix hyphenation
        raw = " ".join(m["texts"]).strip()
        raw = re.sub(r"-\s+\n?", "", raw)  # join broken words
        line_text = re.sub(r"\s+", " ", raw)
        if not line_text:
            continue
        avg_conf = float(np.mean(m["conf"])) if m["conf"] else 1.0
        lines.append({
            "text": line_text,
            "left": float(min(m["xs"])),
            "y": float(np.median(m["ys"])) if m["ys"] else 0.0,
            "conf": avg_conf
        })
    return lines


def lines_to_markdown(lines: List[Dict[str, Any]], conf_threshold: float = 0.50) -> str:
    """
    Convert normalized OCR line dicts to simple Markdown with heading + bullets.
    Uses left offsets to infer indentation.
    """
    if np is None:
        # minimal fallback: no indentation estimation
        keep = [l for l in lines if l.get("conf", 1.0) >= conf_threshold and l.get("text")]
        if not keep:
            return ""
        title = keep[0]["text"]
        md = [f"# {title}"] if (len(title) <= 80 and not title.endswith(".")) else []
        for l in (keep[1:] if md else keep):
            md.append(f"- {l['text']}")
        return "\n".join(md)

    keep = [l for l in lines if l.get("conf", 1.0) >= conf_threshold and (l.get("text") or "").strip()]
    if not keep:
        return ""

    # title heuristic: short, capitalized, non-sentence
    first = keep[0]["text"].strip()
    md: List[str] = []
    if len(first) <= 80 and first[:1].upper() == first[:1] and not first.endswith("."):
        md.append("# " + first)
        keep = keep[1:]

    # indentation from left offsets
    lefts = np.array([l["left"] for l in keep], dtype=float)
    base = float(np.percentile(lefts, 10))
    spread = lefts - base
    indent_unit = max(12.0, float(np.percentile(spread[spread > 0], 50)) if np.any(spread > 0) else 24.0)

    # dedupe consecutive duplicates
    last_txt = None
    for l in keep:
        txt = re.sub(r"\s+", " ", l["text"].strip())
        if not txt or txt == last_txt:
            continue
        last_txt = txt

        # preserve numbering if present
        num_prefix = re.match(r"^(\d{1,3}[.)]|[A-Z]\.)\s+", txt)
        depth = int(max(0, round((l["left"] - base) / indent_unit)))
        bullet = ("  " * depth) + ("- " if not num_prefix else "")
        md.append(bullet + txt)

    return "\n".join(md)


def ocr_slide_to_markdown(
    image_path: str,
    ocr_engine,
    use_paddle: bool = False,
    max_width: int = 1600
) -> str:
    """
    OCR a rendered slide image into structured Markdown.
    """
    if Image is None:
        raise RuntimeError("Pillow is required for OCR processing (PIL.Image).")
    if np is None:
        raise RuntimeError("numpy is required for OCR processing.")

    img = Image.open(image_path).convert("RGB")
    if img.width > max_width:
        ratio = max_width / img.width
        img = img.resize((max_width, int(img.height * ratio)))

    img_np = np.array(img)

    # optional pre-processing for sharper OCR
    try:
        img_np = _preprocess_for_ocr(img_np, enable_denoise=True, enable_binarize=True)
    except Exception:
        pass

    items = ocr_image_paddle(ocr_engine, img_np) if use_paddle else ocr_image_easy(ocr_engine, img_np)
    lines = normalize_lines(items)
    return lines_to_markdown(lines)


# -----------------------------
# End-to-end: slides → text
# -----------------------------
def slides_to_clean_text(
    pptx_path: str,
    out_dir: str = "slide_text_output",
    ocr_engine_name: str = "easyocr",
    lang: str = "en",
    prefer_windows_com: bool = True,
    dpi: int = 220
) -> Dict[str, Any]:
    """
    Render PPTX/PDF to images, OCR to Markdown per slide, capture notes/alt text,
    and zip all outputs.

    Returns:
        {
            "images": [paths],
            "per_slide_markdown": [paths],
            "concat_txt": path,
            "zip": path,
            "out_dir": path
        }
    """
    if np is None or Image is None:
        raise RuntimeError("numpy and Pillow are required for slide text extraction.")

    p = Path(pptx_path)
    out = Path(out_dir); out.mkdir(parents=True, exist_ok=True)
    slides_dir = out / "slides"; slides_dir.mkdir(exist_ok=True)

    ext = p.suffix.lower()
    slide_imgs: List[Path] = []

    # Render images
    if ext == ".pdf":
        # PDF via PyMuPDF
        try:
            import fitz  # PyMuPDF
        except Exception as e:
            raise RuntimeError("PyMuPDF (pymupdf) is required for PDF input.") from e

        doc = fitz.open(str(p))
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            img_path = slides_dir / f"slide_{i:03d}.png"
            pix.save(str(img_path))
            slide_imgs.append(img_path)
    else:
        # PPTX via render helper (COM or fallback implemented there)
        if render_slides_to_images is None:
            raise RuntimeError("render_slides_to_images is not available for PPTX rendering.")
        slide_imgs = [Path(s) for s in render_slides_to_images(
            str(p), out_dir=str(slides_dir), dpi=dpi, prefer_windows_com=prefer_windows_com
        )]

    slide_imgs = sorted(slide_imgs)

    # Notes/alt only for PPTX
    meta = extract_notes_and_alttext(str(p)) if ext == ".pptx" else {}

    # OCR engine
    if get_ocr_engine is None:
        raise RuntimeError("get_ocr_engine is not available; OCR module missing.")
    ocr_engine, use_paddle = get_ocr_engine(ocr_engine_name, lang)

    # OCR each slide → Markdown
    per_slide_md_paths: List[Path] = []
    all_md_blocks: List[str] = []
    try:
        for idx, img_path in enumerate(tqdm(slide_imgs, desc="OCR slides"), 1):
            md = ocr_slide_to_markdown(str(img_path), ocr_engine, use_paddle=use_paddle)
            notes = meta.get(idx, {}).get("notes", "").strip()
            alt_texts = meta.get(idx, {}).get("alt_texts", [])

            if notes:
                md += f"\n\n> Notes:\n> {notes}"
            if alt_texts:
                md += "\n\n> Alt text:\n" + "\n".join([f"> - {t}" for t in alt_texts])

            if not md.strip():
                md = f"# Slide {idx}\n- (No text detected)"

            slide_md_path = out / f"slide_{idx:03d}.md"
            slide_md_path.write_text(md, encoding="utf-8")

            per_slide_md_paths.append(slide_md_path)
            all_md_blocks.append(f"<!-- SLIDE {idx} -->\n{md}\n")
    finally:
        cleanup_ocr_cache(ocr_engine)

    # Concatenate
    concat_path = out / "slides_concatenated.txt"
    concat_path.write_text("\n\n".join(all_md_blocks), encoding="utf-8")

    # Zip artifacts
    zip_path = out / "slide_text_outputs.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.write(concat_path, arcname=concat_path.name)
        for mdp in per_slide_md_paths:
            zf.write(mdp, arcname=mdp.name)
        for img in slide_imgs:
            zf.write(img, arcname=f"slides/{Path(img).name}")

    return {
        "images": [str(p) for p in slide_imgs],
        "per_slide_markdown": [str(p) for p in per_slide_md_paths],
        "concat_txt": str(concat_path),
        "zip": str(zip_path),
        "out_dir": str(out)
    }
