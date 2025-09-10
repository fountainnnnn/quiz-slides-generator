# -*- coding: utf-8 -*-
from typing import Any, Dict, List, Optional
from pathlib import Path
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

ACCENT_Q = RGBColor(37, 99, 235)    # blue
ACCENT_A = RGBColor(16, 185, 129)   # green
ACCENT_X = RGBColor(234, 179, 8)    # amber
FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

# ---------------------------
# Low-level text helpers
# ---------------------------
def _add_run(p, text: str, *, bold: bool = False, mono: bool = False, size: int = 24,
             color: Optional[RGBColor] = None):
    r = p.add_run()
    r.text = text
    r.font.name = FONT_MONO if mono else FONT_SANS
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return r

def _set_textframe(tf, text: str, *, size: int = 24, mono: bool = False, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text.splitlines()):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        _add_run(p, line, mono=mono, size=size)

def _set_title(tf, title: str, *, color: Optional[RGBColor] = None):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, title, bold=True, size=30, color=color)

def _add_footer(slide, text: str):
    left, top, width, height = Inches(8.5), Inches(7.0), Inches(1.3), Inches(0.4)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _add_run(p, text, size=12)

def _add_thumbnail(slide, thumbnail_path: Optional[str], *, max_w=3.2, top=1.2):
    if not thumbnail_path:
        return
    p = Path(thumbnail_path)
    if not p.exists():
        return
    try:
        # right-side thumbnail, fixed width
        slide.shapes.add_picture(str(p), Inches(10 - max_w - 0.6), Inches(top), width=Inches(max_w))
    except Exception:
        pass

# ---------------------------
# Slide builders
# ---------------------------
def _add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    s0 = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    try:
        _set_textframe(s0.shapes.title.text_frame, title, size=44)
        if len(s0.placeholders) > 1:
            sub = subtitle or time.strftime("Generated on %Y-%m-%d %H:%M")
            _set_textframe(s0.placeholders[1].text_frame, sub, size=18)
    except Exception:
        # Fallback to textboxes if placeholders differ
        tf = s0.shapes.add_textbox(Inches(1), Inches(1.0), Inches(8), Inches(1.5)).text_frame
        _set_textframe(tf, title, size=44)
        tf2 = s0.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(0.8)).text_frame
        _set_textframe(tf2, subtitle or time.strftime("Generated on %Y-%m-%d %H:%M"), size=18)

def _add_q_slide(prs: Presentation, idx: int, total: int, src: int, body: str,
                 *, thumbnail: Optional[str] = None, mono: bool = False):
    s = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    # Title
    try:
        _set_title(s.shapes.title.text_frame, f"Question {idx}/{total} — Slide {src}", color=ACCENT_Q)
    except Exception:
        t = s.shapes.add_textbox(Inches(1), Inches(0.9), Inches(8), Inches(0.8)).text_frame
        _set_textframe(t, f"Question {idx}/{total} — Slide {src}", size=30)
    # Content
    tf = s.placeholders[1].text_frame if len(s.placeholders) > 1 else s.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(7.8), Inches(4.8)
    ).text_frame
    _set_textframe(tf, body, size=28 if not mono else 26, mono=mono)
    # Thumbnail (optional)
    _add_thumbnail(s, thumbnail)
    _add_footer(s, f"{idx* (2) - 1 if True else idx}/{total}")  # simple counter; adjusted below in loop

def _add_answer_slide(prs: Presentation, idx: int, total: int, answer: str, explanation: str):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    # Title
    try:
        _set_title(s.shapes.title.text_frame, f"Answer {idx}/{total}", color=ACCENT_A)
    except Exception:
        t = s.shapes.add_textbox(Inches(1), Inches(0.9), Inches(8), Inches(0.8)).text_frame
        _set_textframe(t, f"Answer {idx}/{total}", size=30)
    # Content: bold "Answer:" then explanation as new paragraphs
    tf = s.placeholders[1].text_frame if len(s.placeholders) > 1 else s.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(7.8), Inches(4.8)
    ).text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    # Answer line
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, "Answer: ", bold=True, size=28, color=ACCENT_A)
    _add_run(p, answer, size=28)
    # Explanation (new paragraph if present)
    if explanation and explanation.strip():
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        _add_run(p2, explanation.strip(), size=24)
    _add_footer(s, f"{idx* (2)}/{total}")  # adjusted below in loop
    return s

def _add_explanation_slide(prs: Presentation, idx: int, total: int, explanation: str):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    # Title
    try:
        _set_title(s.shapes.title.text_frame, f"Why? {idx}/{total}", color=ACCENT_X)
    except Exception:
        t = s.shapes.add_textbox(Inches(1), Inches(0.9), Inches(8), Inches(0.8)).text_frame
        _set_textframe(t, f"Why? {idx}/{total}", size=30)
    # Content
    tf = s.placeholders[1].text_frame if len(s.placeholders) > 1 else s.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(7.8), Inches(4.8)
    ).text_frame
    _set_textframe(tf, explanation.strip(), size=24)
    return s

# ---------------------------
# Public API
# ---------------------------
def build_qa_deck(
    qa: List[Dict[str, Any]],
    slide_images: List[str],
    out_path: str,
    deck_title: str = "Auto Q&A Deck",
    include_thumbnails: bool = False,
    include_explanations: bool = True,
    source_name: str = ""
):
    """
    qa: list of {type, question, options, answer, explanation, source_slide_index, code?}
    slide_images: list of slide image paths (for thumbnails)
    """
    prs = Presentation()
    subtitle = time.strftime("Generated on %Y-%m-%d %H:%M") + (f" • from {source_name}" if source_name else "")
    _add_title_slide(prs, deck_title, subtitle)

    total = len(qa)
    for i, q in enumerate(qa, 1):
        typ = str(q.get("type", "")).lower()
        src = int(q.get("source_slide_index", 1) or 1)
        thumb = slide_images[src - 1] if include_thumbnails and 0 <= src - 1 < len(slide_images) else None

        # Build question body
        if typ == "mcq":
            opts = [str(o) for o in (q.get("options") or []) if str(o).strip()]
            opts_block = "\n".join([f"{chr(65 + j)}. {opt}" for j, opt in enumerate(opts[:4])])
            q_body = (q.get("question") or "").strip()
            body = f"{q_body}\n\n{opts_block}\n\n(Select one)"
            _add_q_slide(prs, i, total, src, body, thumbnail=thumb, mono=False)

            ans = (q.get("answer") or "").strip()
            exp = (q.get("explanation") or "").strip()
            _add_answer_slide(prs, i, total, ans, exp)
            if include_explanations and exp:
                _add_explanation_slide(prs, i, total, exp)

        elif typ == "code_fill":
            q_body = (q.get("question") or "").strip()
            code = (q.get("code") or "").strip()
            if code:
                q_body = f"{q_body}\n\n{code}"
            # show options if present
            opts = [str(o) for o in (q.get("options") or []) if str(o).strip()]
            if opts:
                q_body += "\n\nOptions: " + ", ".join(opts[:6])
            _add_q_slide(prs, i, total, src, q_body, thumbnail=thumb, mono=True)

            ans = (q.get("answer") or "").strip()
            exp = (q.get("explanation") or "").strip()
            _add_answer_slide(prs, i, total, ans, exp)
            if include_explanations and exp:
                _add_explanation_slide(prs, i, total, exp)

        else:  # theory / default
            q_body = (q.get("question") or "").strip()
            _add_q_slide(prs, i, total, src, q_body, thumbnail=thumb, mono=False)

            ans = (q.get("answer") or "").strip()
            exp = (q.get("explanation") or "").strip()
            _add_answer_slide(prs, i, total, ans, exp)
            if include_explanations and exp:
                _add_explanation_slide(prs, i, total, exp)

    prs.save(out_path)
    return out_path
